# app.py - VERSÃO COM SUPORTE A DOCX E PPTX

import os
import io
import csv
import json
import fitz  # PyMuPDF
import requests
from bs4 import BeautifulSoup
from flask import (
    Flask, request, render_template, make_response, send_file, session, redirect, url_for, g
)
from flask_babel import Babel, _, lazy_gettext as _l
from werkzeug.utils import secure_filename
from urllib.parse import urlparse
from dotenv import load_dotenv

# --- NOVAS Importações ---
try:
    import docx
except ImportError:
    docx = None
    print("AVISO: Biblioteca 'python-docx' não encontrada. Extração de .docx desabilitada.")
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    print("AVISO: Biblioteca 'python-pptx' não encontrada. Extração de .pptx desabilitada.")
# --- Fim Novas Importações ---

load_dotenv()

# --- Configuração do App e Babel ---
app = Flask(__name__)
app.secret_key = os.environ.get('FLASK_SECRET_KEY', 'temp-unsafe-key-for-initial-run-only-define-env-var') # Definir valor padrão
if app.secret_key == 'temp-unsafe-key-for-initial-run-only-define-env-var':
     print("\n*** ATENÇÃO: Usando chave secreta temporária e insegura. Defina FLASK_SECRET_KEY no seu ambiente! ***\n")

app.config['LANGUAGES'] = {'pt': 'Português', 'en': 'English', 'es': 'Español'}
app.config['BABEL_DEFAULT_LOCALE'] = 'pt'

def get_locale():
    lang = session.get('language')
    if lang in app.config['LANGUAGES']:
        return lang
    return request.accept_languages.best_match(app.config['LANGUAGES'].keys()) or app.config['BABEL_DEFAULT_LOCALE']

babel = Babel(app, locale_selector=get_locale)

@app.route('/language/<lang>')
def set_language(lang=None):
    if lang in app.config['LANGUAGES'].keys():
        session['language'] = lang
    return redirect(request.referrer or url_for('index'))

@app.context_processor
def inject_conf_var():
    lang = session.get('language', get_locale())
    return dict(CURRENT_LANG=lang)
# --- Fim Configuração / Idioma ---


# --- Constantes e Funções Utilitárias ---
# MODIFICADO: Adicionar novas extensões permitidas
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx'}
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024 # 16 MB
REQUESTS_TIMEOUT = 20

def allowed_file(filename):
    """Verifica se a extensão do arquivo é permitida."""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def is_valid_url(url):
    """Validação básica de formato de URL."""
    try:
        result = urlparse(url)
        return all([result.scheme in ['http', 'https'], result.netloc])
    except ValueError:
        return False
# --- Fim Constantes / Utilitários ---


# --- Funções de Extração (com mensagens traduzíveis) ---
def extract_pdf_text(byte_stream, source_description="PDF"):
    """Extrai texto de um stream de bytes de um PDF."""
    text = ""
    doc = None
    try:
        doc = fitz.open(stream=byte_stream, filetype="pdf")
        if not doc.is_pdf:
             raise ValueError(_('The content from "%(source)s" does not appear to be a valid PDF.', source=source_description))
        if doc.is_encrypted:
             if not doc.authenticate(""):
                  raise ValueError(_('The PDF from "%(source)s" is password protected and cannot be processed.', source=source_description))
        for page_num in range(len(doc)):
            try:
                page = doc.load_page(page_num)
                text += page.get_text("text") + "\n" # Adiciona nova linha entre páginas
            except Exception as page_error:
                print(f"Aviso: Falha ao processar página {page_num + 1} de '{source_description}': {page_error}")
                text += _('[Error processing page %(num)s]', num=page_num + 1) + "\n"
        if not text.strip():
             raise ValueError(_('No extractable text found in the PDF "%(source)s". It might be an image file, empty, or contain only graphics.', source=source_description))
        return text.strip()
    except fitz.fitz.FileDataError as e:
        raise ValueError(_('Invalid or corrupted PDF file from "%(source)s": %(error)s', source=source_description, error=str(e)))
    except ValueError as ve:
        raise ve # Re-lança erros específicos já traduzidos
    except Exception as e:
        print(f"Erro inesperado ao ler PDF de {source_description}: {type(e).__name__} - {e}")
        if "cannot open" in str(e) or "syntax error" in str(e):
             raise ValueError(_('Invalid or unsupported PDF format from "%(source)s": %(error)s', source=source_description, error=str(e)))
        else:
             raise ValueError(_('Unexpected error processing the PDF from "%(source)s": %(error)s', source=source_description, error=str(e)))
    finally:
         if doc: doc.close()

def extract_html_text(html_content, source_url="Página Web"):
    """Extrai texto visível de conteúdo HTML."""
    try:
        soup = BeautifulSoup(html_content, 'lxml')
        for element in soup(["script", "style", "head", "title", "meta", "[document]", "header", "footer", "nav", "aside", "form", "button", "select", "textarea", "img", "svg", "iframe", "noscript"]):
            element.decompose()
        text = soup.get_text(separator='\n', strip=True) # Usar nova linha como separador
        if not text:
            raise ValueError(_('No main text found at "%(url)s". The page might be empty, heavily reliant on JavaScript, or the main content was removed during parsing.', url=source_url))
        return text
    except Exception as e:
        print(f"Erro ao processar HTML de {source_url}: {type(e).__name__} - {e}")
        raise ValueError(_('Error parsing HTML content from URL "%(url)s": %(error)s', url=source_url, error=str(e)))

# --- NOVA Função: Extrair texto de DOCX ---
def extract_docx_text(byte_stream, source_description="Word Document"):
    """Extrai texto de um stream de bytes de um DOCX."""
    if not docx:
        raise RuntimeError(_("Word document processing library (python-docx) is not installed."))
    text = ""
    try:
        document = docx.Document(byte_stream)
        for para in document.paragraphs:
            text += para.text + "\n"
        # Poderia adicionar extração de tabelas aqui se necessário
        # for table in document.tables:
        #     for row in table.rows:
        #         for cell in row.cells:
        #             text += cell.text + "\t" # Tab para separar células
        #         text += "\n" # Nova linha para cada linha da tabela
        if not text.strip():
            raise ValueError(_('No text found in the Word document "%(source)s".', source=source_description))
        return text.strip()
    except Exception as e: # Captura exceções genéricas e específicas do docx se conhecidas
        print(f"Erro ao processar DOCX de {source_description}: {type(e).__name__} - {e}")
        # Tenta identificar erro comum de formato inválido
        if "File is not a zip file" in str(e) or isinstance(e, docx.opc.exceptions.PackageNotFoundError):
            raise ValueError(_('Invalid or corrupted Word document format (DOCX expected) from "%(source)s".', source=source_description))
        else:
            raise ValueError(_('Error processing Word document "%(source)s": %(error)s', source=source_description, error=str(e)))

# --- NOVA Função: Extrair texto de PPTX ---
def extract_pptx_text(byte_stream, source_description="PowerPoint Presentation"):
    """Extrai texto de um stream de bytes de um PPTX."""
    if not Presentation:
         raise RuntimeError(_("PowerPoint processing library (python-pptx) is not installed."))
    text = ""
    try:
        prs = Presentation(byte_stream)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            # Texto das formas no slide
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)
            # Texto das notas (se existir)
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                   slide_text.append(notes_frame.text) # Adiciona todo o texto das notas

            if slide_text:
                 text += _("--- Slide %(num)s ---", num=i+1) + "\n" + "\n".join(slide_text) + "\n\n"

        if not text.strip():
            raise ValueError(_('No text found in the PowerPoint presentation "%(source)s".', source=source_description))
        return text.strip()
    except Exception as e: # Captura exceções genéricas e específicas do pptx se conhecidas
        print(f"Erro ao processar PPTX de {source_description}: {type(e).__name__} - {e}")
        # Tenta identificar erro comum de formato inválido
        if "File is not a zip file" in str(e) or isinstance(e, pptx.exc.PackageNotFoundError):
             raise ValueError(_('Invalid or corrupted PowerPoint presentation format (PPTX expected) from "%(source)s".', source=source_description))
        else:
            raise ValueError(_('Error processing PowerPoint presentation "%(source)s": %(error)s', source=source_description, error=str(e)))
# --- Fim Funções de Extração ---


# --- Rotas Principais ---
@app.route('/')
def index():
    """Renderiza a página inicial."""
    # MODIFICADO: Passa as extensões permitidas para o template
    allowed_ext_string = ", ".join([f".{ext.upper()}" for ext in ALLOWED_EXTENSIONS])
    return render_template('index.html', allowed_extensions=allowed_ext_string)

@app.route('/process', methods=['POST'])
def process_input():
    """Processa upload de PDF/DOCX/PPTX OU URL (PDF/HTML/DOCX/PPTX)."""
    file_input = request.files.get('file_input') # Nome do campo de arquivo genérico
    input_url = request.form.get('input_url', '').strip()
    source_description = _("input")
    byte_stream = None
    filename_base = "extracted_data"
    extracted_text = None
    file_extension = None

    try:
        # --- Upload ---
        if file_input and file_input.filename != '':
            filename = secure_filename(file_input.filename)
            if not allowed_file(filename):
                allowed_types = ", ".join(ALLOWED_EXTENSIONS)
                raise ValueError(
                    _('Invalid file type: "%(filename)s". Please upload a supported file type: %(types)s.',
                      filename=filename, types=allowed_types)
                )

            source_description = filename
            filename_base = filename.rsplit('.', 1)[0] if '.' in filename else filename
            file_extension = filename.rsplit('.', 1)[1].lower()

            byte_stream = io.BytesIO()
            file_input.save(byte_stream)
            byte_stream.seek(0)

            # MODIFICADO: Chama a função de extração correta baseada na extensão
            if file_extension == 'pdf':
                extracted_text = extract_pdf_text(byte_stream, source_description)
            elif file_extension == 'docx':
                extracted_text = extract_docx_text(byte_stream, source_description)
            elif file_extension == 'pptx':
                extracted_text = extract_pptx_text(byte_stream, source_description)
            else:
                # Isso não deveria acontecer devido à verificação allowed_file, mas é uma segurança
                raise ValueError(_("Internal error: File type '%(ext)s' was allowed but not handled.", ext=file_extension))

        # --- URL ---
        elif input_url:
            if not is_valid_url(input_url):
                raise ValueError(_('Invalid URL format.'))
            source_description = input_url

            # Define o nome base do arquivo a partir da URL
            try:
                 parsed_url = urlparse(input_url)
                 url_filename = os.path.basename(parsed_url.path); fn = url_filename
                 if fn and '.' in fn: filename_base = fn.rsplit('.', 1)[0]
                 elif fn: filename_base = fn
                 else: filename_base = parsed_url.netloc.replace('.', '_') or "extracted_from_url"
            except Exception: filename_base = "extracted_from_url"

            # Tenta obter o arquivo da URL
            try:
                headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
                response = requests.get(input_url, timeout=REQUESTS_TIMEOUT, headers=headers, allow_redirects=True, stream=True) # stream=True é importante para grandes arquivos
                response.raise_for_status()

                content_type = response.headers.get('content-type', '').lower()
                content_length = int(response.headers.get('content-length', 0))

                if content_length > app.config['MAX_CONTENT_LENGTH']:
                     raise ValueError(_('The content at the URL is too large (larger than %(max_size)sMB).', max_size=app.config['MAX_CONTENT_LENGTH']//(1024*1024)))

                # Determina o tipo de arquivo pela URL ou Content-Type
                url_path_lower = urlparse(input_url).path.lower()
                is_pdf = 'application/pdf' in content_type or url_path_lower.endswith('.pdf')
                is_docx = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type or url_path_lower.endswith('.docx')
                is_pptx = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' in content_type or url_path_lower.endswith('.pptx')
                is_html = 'text/html' in content_type or 'application/xhtml+xml' in content_type

                # Lê o conteúdo em um stream de bytes
                byte_stream_url = io.BytesIO()
                read_bytes = 0
                for chunk in response.iter_content(chunk_size=8192):
                    read_bytes += len(chunk)
                    if read_bytes > app.config['MAX_CONTENT_LENGTH']:
                         byte_stream_url.close()
                         raise ValueError(_('The content at the URL is too large (larger than %(max_size)sMB).', max_size=app.config['MAX_CONTENT_LENGTH']//(1024*1024)))
                    byte_stream_url.write(chunk)
                byte_stream_url.seek(0)

                # MODIFICADO: Chama a função correta baseada no tipo detectado
                if is_pdf:
                    extracted_text = extract_pdf_text(byte_stream_url, source_description)
                elif is_docx:
                    extracted_text = extract_docx_text(byte_stream_url, source_description)
                elif is_pptx:
                    extracted_text = extract_pptx_text(byte_stream_url, source_description)
                elif is_html and not (is_pdf or is_docx or is_pptx): # HTML como fallback se não for outro tipo conhecido
                    # Para HTML, precisamos do conteúdo como string, então lemos do stream de bytes
                    byte_stream_url.seek(0)
                    # Tenta decodificar usando a codificação da resposta ou UTF-8 como fallback
                    encoding = response.encoding or 'utf-8'
                    html_content = byte_stream_url.read().decode(encoding, errors='replace')
                    extracted_text = extract_html_text(html_content, source_description)
                else:
                    # Fecha o stream antes de lançar o erro
                    byte_stream_url.close()
                    supported_types = "PDF, DOCX, PPTX, HTML"
                    raise ValueError(_('The URL does not point to a supported file type (%(types)s). Detected Content-Type: %(ctype)s', types=supported_types, ctype=content_type or 'N/A'))

                # Fecha o stream da URL após o processamento
                byte_stream_url.close()

            except requests.exceptions.Timeout:
                raise ValueError(_('Timeout (%(sec)ss) when trying to access the URL: %(url)s', sec=REQUESTS_TIMEOUT, url=input_url))
            except requests.exceptions.RequestException as e:
                 error_msg = str(e).split('\n')[0] # Mensagem de erro mais curta
                 raise ValueError(_('Network error accessing the URL %(url)s: %(error)s', url=input_url, error=error_msg))

        # --- Nenhum Input ---
        else:
            raise ValueError(_('No file was uploaded or valid URL provided.'))

        # --- Resultado ---
        if extracted_text:
            session['extracted_text'] = extracted_text
            session['original_filename'] = filename_base
            # Renderiza a página de resultados com o texto e a origem
            return render_template('results.html', text=extracted_text, source=source_description)
        else:
             # Este caso pode ocorrer se a extração não retornar texto mas não lançar erro (improvável com as verificações atuais)
             raise ValueError(_('Could not extract text from "%(source)s" (empty result after processing).', source=source_description))

    except ValueError as e: # Erros esperados e traduzidos
         return render_template('results.html', error=str(e), source=source_description) # Passa a origem para contexto
    except RuntimeError as e: # Erros como bibliotecas faltando
         return render_template('results.html', error=str(e), source=source_description)
    except Exception as e: # Erros inesperados
        print(f"ERRO INESPERADO no processamento: {type(e).__name__} - {e}")
        import traceback
        traceback.print_exc() # Log completo do erro no console do servidor
        return render_template('results.html', error=_('An unexpected error occurred during processing. Check server logs for details.'), source=source_description)
    finally:
         # Garante que o stream de upload seja fechado se foi aberto
         if byte_stream and not byte_stream.closed:
            byte_stream.close()


@app.route('/download/<format>')
def download_file(format):
    """Fornece o texto extraído para download no formato solicitado."""
    extracted_text = session.get('extracted_text')
    base_filename = session.get('original_filename', 'extracted_data')

    if not extracted_text:
        # Talvez redirecionar para a página inicial com uma mensagem flash seria melhor?
        return _('Error: No extracted text found in the current session for download. Please try processing a file or URL again.'), 404

    allowed_formats = {'txt', 'csv', 'json'}
    if format not in allowed_formats: return _('Invalid download format specified.'), 400
    download_filename = f"{base_filename}.{format}"

    buffer = None
    mimetype = 'application/octet-stream'
    try:
        if format == 'txt':
            buffer = io.BytesIO(extracted_text.encode('utf-8'))
            mimetype = 'text/plain; charset=utf-8'
        elif format == 'csv':
            string_buffer = io.StringIO()
            # Usando QUOTE_MINIMAL para evitar aspas desnecessárias, mas QUOTE_NONNUMERIC ou ALL pode ser mais seguro
            writer = csv.writer(string_buffer, quoting=csv.QUOTE_MINIMAL)
            writer.writerow([_('Extracted Text')]) # Cabeçalho traduzido
            # Divide o texto em linhas e escreve cada linha como uma linha CSV
            # Isso assume que o texto extraído não contém vírgulas que precisam ser escapadas de forma complexa.
            # Para dados mais complexos, pode ser necessário um tratamento mais robusto.
            lines = extracted_text.splitlines()
            for line in lines:
                writer.writerow([line]) # Escreve cada linha como uma única coluna
            string_buffer.seek(0)
            buffer = io.BytesIO(string_buffer.getvalue().encode('utf-8-sig')) # utf-8-sig para melhor compatibilidade com Excel
            string_buffer.close()
            mimetype = 'text/csv; charset=utf-8'
        elif format == 'json':
            # Estrutura JSON simples: um objeto com uma chave "lines" contendo uma lista de strings (linhas)
            lines = extracted_text.splitlines()
            data = {
                "source": session.get('original_filename', 'unknown'),
                "extracted_lines": lines
                }
            json_str = json.dumps(data, ensure_ascii=False, indent=2) # ensure_ascii=False para caracteres não-latinos
            buffer = io.BytesIO(json_str.encode('utf-8'))
            mimetype = 'application/json; charset=utf-8'

        # Envia o arquivo
        return send_file(
            buffer,
            mimetype=mimetype,
            as_attachment=True,
            download_name=download_filename
        )

    except Exception as e:
        print(f"ERRO DOWNLOAD: Falha ao gerar/enviar arquivo {format}: {type(e).__name__} - {e}")
        # Fecha o buffer se ele foi criado e ainda está aberto
        if buffer and hasattr(buffer, 'closed') and not buffer.closed:
            buffer.close()
        # Retorna uma mensagem de erro genérica para o usuário
        return _('An error occurred while generating the file for download.'), 500
    # O 'finally' não é estritamente necessário aqui pois o 'send_file' do Flask fecha o buffer após o envio.

# --- Fim Rotas ---

# Roda o servidor
if __name__ == '__main__':
    # Lembre-se de definir debug=False para produção e usar um servidor WSGI (como Gunicorn ou Waitress)
    # Exemplo: gunicorn -w 4 app:app
    # O host 0.0.0.0 torna acessível na rede local. Mude para 127.0.0.1 se quiser apenas local.
    print(f"Flask app running with ALLOWED_EXTENSIONS: {ALLOWED_EXTENSIONS}")
    app.run(debug=True, host='0.0.0.0', port=5000)